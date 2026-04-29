import os,fitz,requests,re,math,docx,json,io,pytesseract,asyncio
from typing import Optional
from dotenv import load_dotenv
from sqlalchemy.orm import Session
from passlib.context import CryptContext
from app.models.rfp_models import User,KeystoneFile
from app.config import client, index
from pptx import Presentation
from PyPDF2 import PdfReader
from fastapi import HTTPException
from app.models import * 
import pandas as pd
from PIL import Image
from sqlalchemy import select

load_dotenv()
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")

from app.core.llm_client import get_llm_client
from app.core.prompts import (question_prompt,
                              mode_block_prompt,
                              answer_generation_prompt, 
                              summary_and_analysis_prompt, 
                              summary_format_prompt, 
                              search_queries_prompt, 
                              questions_grouped_prompt,
                              generate_score_prompt, 
                              classification_prompt)
from app.core.llm_client.openai import OpenAIEmbeddingClient
from sqlalchemy.ext.asyncio import AsyncSession



async def _complete_with_fallback(
    provider: str,
    prompt: str,
    system_prompt: Optional[str] = None,
    fallback_providers: list[str] = None,
) -> str:
    """
    Try primary model first, then fallback models in order.
    """
    all_providers = [provider] + (fallback_providers or ["gpt-4o-mini", "gpt-5.4"])
    last_exception = None

    for current_provider in all_providers:
        try:
            client = get_llm_client(current_provider)

            if system_prompt is None:
                return await client.complete(prompt=prompt)
            else:
                return await client.complete(
                    prompt=prompt,
                    system=system_prompt
                )

        except Exception as e:
            last_exception = e
            print(f"[WARNING] Provider '{current_provider}' failed: {e}. Trying next...")
            continue

    raise RuntimeError(
        f"All providers failed. Last error: {last_exception}"
    ) from last_exception


def extract_text_from_pdf(pdf_file: bytes) -> str:
    text = ""
    with fitz.open(stream=pdf_file, filetype="pdf") as doc:
        for page in doc:
            extracted_text = page.get_text()

            if extracted_text and extracted_text.strip():
                text += extracted_text
            else:
                pix = page.get_pixmap(dpi=300)
                img_data = pix.tobytes("png")
                image = Image.open(io.BytesIO(img_data))
                ocr_text = pytesseract.image_to_string(image)
                text += ocr_text

    if not text.strip():
        raise ValueError("Unable to extract text from PDF: possibly fully scanned or corrupted.")

    return text


async def generate_search_queries(
    rfp_text: str,
    provider: str,
    fallback_providers: list[str] = None
) -> list:
    """
    Generate exactly 12 highly targeted Google search queries based on RFP text.
    """
    prompt = search_queries_prompt(rfp_text)
    system_prompt = "You generate Google search queries to build complete company profiles from RFPs."

    content = await _complete_with_fallback(provider, prompt, system_prompt, fallback_providers)
    return [line.strip(" -•") for line in content.split("\n") if line.strip()]

async def extract_company_background_from_rfp(
    rfp_text: str,
    provider: str,
    fallback_providers: list[str] = None
) -> str:
    """
    Extracts 3 fully detailed sections from an RFP:
    1. Purpose of the RFP (including Scope of Work, Buyer Priorities & Win Themes)
    2. Company Background
    3. Submission Details & Requirements
    """

    prompt = summary_and_analysis_prompt(rfp_text)

    system_prompt = (
        "You are a meticulous RFP extraction specialist with perfect attention to detail. "
        "Your extractions are comprehensive, accurate, and complete. You NEVER add information "
        "not present in the source document. You NEVER miss important details. You read entire "
        "documents systematically and extract every relevant piece of information. "
        "You ALWAYS decompose Scope of Work items individually — never merging them — and you "
        "ONLY pull scope items from the section explicitly labeled Scope of Work, never from "
        "Priorities, Goals, or Focus Areas sections. "
        "Your Section 3 extractions are especially thorough, capturing every single submission "
        "requirement. You work methodically through checklists to ensure nothing is overlooked."
    )
    content = await _complete_with_fallback(provider, prompt, system_prompt, fallback_providers)

    return content

async def questions_grouped_function(
    rfp_text: str,
    custom_instruction: str,
    provider: str,
    fallback_providers: list[str] = None
) -> dict:
    """
    Generate high-quality proposal questions grouped by RFP section, based on:
    1. The full RFP text
    2. Custom instructions from the admin (if provided)

    CUSTOM INSTRUCTION GUIDELINES:
    - If admin provides specific instructions, use them to guide question generation.
    - Tailor questions to align with admin's focus areas and priorities.
    - If admin narrows scope, exclude irrelevant questions and focus on specified topics."""
    
    prompt = questions_grouped_prompt(rfp_text, custom_instruction)
    system_prompt = "Return ONLY valid JSON."

    content = await _complete_with_fallback(provider, prompt, system_prompt, fallback_providers)

    content = content.strip().replace("```json", "").replace("```", "").strip()
    try:
        return json.loads(content)
    except Exception:
        raise HTTPException(
            status_code=500,
            detail="AI returned invalid JSON in custom question generation"
        )

async def summarize_results_with_llm(
    all_snippets: list,
    rfp_company_text: str,
    provider: str,
    fallback_providers: list[str] = None,
) -> str:

    combined_snippets = "\n".join(all_snippets)
    system_prompt = (
        "You are a meticulous senior RFP analyst and strategy consultant who produces "
        "structured analysis briefs. You extract and organize information with perfect accuracy, "
        "never adding content not in the sources."
    )

    async def generate_section(section_num):
        prompt = summary_format_prompt(section_num, rfp_company_text, combined_snippets)
        return await _complete_with_fallback(
            provider,
            prompt,
            system_prompt,
            fallback_providers
        )

    section1, section2, section3, section4 = await asyncio.gather(
        generate_section(1),
        generate_section(2),
        generate_section(3),
        generate_section(4),
    )

    return f"{section1}\n\n---\n\n{section2}\n\n---\n\n{section3}\n\n---\n\n{section4}"

async def extract_questions_with_llm(
    classification_QaI_results: str,
    provider: str,
    fallback_providers: list[str] = None
):
    
    prompt = question_prompt(classification_QaI_results)
    system_prompt = "Return ONLY strict valid JSON. No markdown. No commentary."
    content = await _complete_with_fallback(provider, prompt, system_prompt, fallback_providers)
    # print("Raw LLM output for question extraction:", content)

    content = content.strip()
    content = content.replace("```json", "").replace("```", "").strip()

    try:
        grouped_questions = json.loads(content)
    except Exception:
        raise HTTPException(status_code=500,
                            detail="AI returned invalid JSON for extracted questions.")

    if not isinstance(grouped_questions, dict):
        raise HTTPException(status_code=500,
                            detail="Questions JSON must be an object.")

    for k, v in grouped_questions.items():
        if not isinstance(v, dict):
            raise HTTPException(status_code=500,
                                detail=f"Invalid section {k}")
        if "section" not in v or "questions" not in v:
            raise HTTPException(status_code=500,
                                detail=f"Missing keys in section {k}")

    return grouped_questions

def get_user_by_email(db: Session, email: str):
    return db.query(User).filter(User.email == email).first()

def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def hash_password(password):
    return pwd_context.hash(password)

async def get_similar_context(question: str, rfp_id: int, top_k: int = 5):
    """
    Retrieve RFP-specific chunks from Pinecone using file_id metadata filter.
    """
    try:
        # embedding = client.embeddings.create(
        #     input=[question],
        #     model="text-embedding-3-small"
        # ).data[0].embedding
        embedding = await OpenAIEmbeddingClient().embed(question)
        
        results = index.query(
            vector=embedding,
            top_k=top_k,
            include_metadata=True,
            filter={"file_id": str(rfp_id)}
        )

        context_texts = [match["metadata"]["text"] for match in results["matches"]]
        sources = [
            {
                "score": match["score"],
                "file_id": match["metadata"].get("file_id"),
                "chunk_index": match["metadata"].get("chunk_index"),
                "snippet": match["metadata"].get("text")[:300],
            }
            for match in results["matches"]
        ]

        return "\n".join(context_texts), sources

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Pinecone retrieval failed: {str(e)}")


def _sanitize_short_name(short_name: str) -> str:
    """
    Validate that short_name is a human-readable client name.
    If it looks like a UUID, hash, or alphanumeric ID, return a safe fallback.
    This is the defensive layer — get_short_name() should already be clean,
    but this ensures nothing bad ever reaches the LLM prompt.
    """
    if not short_name or not short_name.strip():
        return "the City"

    name = short_name.strip()

    # Full UUID: 8-4-4-4-12 hex pattern
    uuid_pattern = re.compile(
        r'^[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}'
        r'-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}$',
        re.IGNORECASE
    )

    # Any string that is purely hex characters and dashes/underscores
    # (catches partial UUIDs like "24ad8e0c" or "24Ad8E0C")
    hex_only_pattern = re.compile(r'^[0-9a-fA-F\-_]{6,}$', re.IGNORECASE)

    # Starts with digits and has no spaces (e.g. "24AA07", "2024RFP")
    leading_digit_code = re.compile(r'^\d+[A-Za-z0-9\-_]*$')

    if (uuid_pattern.match(name)
            or hex_only_pattern.match(name)
            or leading_digit_code.match(name)
            or name.lower() == "the organization"):
        return "the City"

    return name


async def generate_answer_with_context(
    question: str,
    context: str,
    short_name: str,
    existing_answer: str = None,
    edit_instruction: str = None,
    provider: str = "gpt-4o-mini",
    # fallback_providers: list[str] = None
) -> str:
    """
    Generate a new proposal response OR apply a targeted edit to an existing one.

    Parameters
    ----------
    question          : The RFP question being answered.
    context           : RAG-retrieved context (RFP chunks + Keystone company data).
    short_name        : Human-readable client name, e.g. "Duluth" or "McLean".
                        Sanitized internally — UUID/hash values are rejected.
    existing_answer   : Previously generated answer. Required for edit mode.
    edit_instruction  : The specific change the user wants applied. Required for edit mode.
                        Both existing_answer AND edit_instruction must be provided
                        to activate edit mode. If either is missing, generate mode runs.
    """

    # Sanitize client name before it ever touches the prompt
    short_name = _sanitize_short_name(short_name)

    # Determine mode
    is_edit_mode = bool(
        existing_answer and existing_answer.strip()
        and edit_instruction and edit_instruction.strip()
    )
    formatting_signals = []
    if is_edit_mode and existing_answer:
        lines = existing_answer.strip().splitlines()
        has_bullets = any(
            line.strip().startswith(("-", "*", "•")) for line in lines
        )
        has_numbered_list = any(
            line.strip() and line.strip()[0].isdigit() and len(line.strip()) > 1
            and line.strip()[1] in (".", ")")
            for line in lines
        )
        paragraph_count = sum(
            1 for i, line in enumerate(lines)
            if line.strip() == "" and i > 0 and lines[i - 1].strip() != ""
        )

        if has_bullets:
            formatting_signals.append("- The existing answer uses bullet points. PRESERVE this bullet-point formatting throughout the edited version.")
        if has_numbered_list:
            formatting_signals.append("- The existing answer uses a numbered list. PRESERVE this numbered-list formatting throughout the edited version.")
        if paragraph_count >= 2:
            formatting_signals.append(f"- The existing answer has {paragraph_count + 1} paragraphs separated by blank lines. PRESERVE these paragraph breaks throughout the edited version.")

    formatting_carry_forward = (
        "\n        ### FORMATTING CARRY-FORWARD (MANDATORY IN EDIT MODE)\n\n"
        + "\n".join(f"        {s}" for s in formatting_signals)
        + "\n"
    ) if formatting_signals else ""


    mode_block = mode_block_prompt(is_edit_mode, existing_answer, edit_instruction, formatting_carry_forward)
    prompt = answer_generation_prompt(question, context, short_name, mode_block)
    system_prompt=(
                "You are a professional RFP response specialist who strictly follows "
                "behavioral authority, factual accuracy, and style rules. "
                "You NEVER use document IDs, UUIDs, filenames, or alphanumeric reference "
                "codes as client names — only the human-readable name explicitly provided. "
                "In edit mode, you apply ONLY the requested change and preserve all other "
                "existing content word for word. "
                "You NEVER begin any response with 'OK', 'Sure', 'Certainly', 'Absolutely','here is the answer' "
                "or any other affirmation — start directly with the answer. "
                "You NEVER fabricate specific facts, figures, budgets, or case study details "
                "not present in the provided context."
                # "Make the generaed output responce answer in bullet points and every bullet points should have a blank line space between each bullet points "
            )
    try:
        content = await _complete_with_fallback(provider, prompt, system_prompt)
        return content

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM generation failed: {str(e)}")

async def analyze_answer_score_only(
    question_text: str,
    answer_text: str,
    provider: str = "openai",
    # fallback_providers: list[str] = None
) -> float:
    prompt = generate_score_prompt(question_text, answer_text)

    system_prompt = "You are a strict RFP evaluator who gives a numeric score based on how well the answer addresses the question. Return ONLY the numeric score as a float from 0.0 to 10.0, with no explanation or text."
    content = await _complete_with_fallback(provider, prompt, system_prompt)
    score_text = content.strip()
    try:
        return float(score_text)
    except ValueError:
        return None

def parse_rfp_summary(summary_text: str):
    """
    Convert LLM markdown summary into structured JSON:
    {
      "purpose": "...",
      "company_background": "...",
      "submission_details": {
          "text": "...",
          "bullets": ["...", "...", ...]
      }
    }
    """
    summary_text = re.sub(r"^#+\s*RFP Analysis Brief\s*", "", summary_text, flags=re.IGNORECASE)

    sections = re.split(r"####\s*Section\s*\d:\s*", summary_text)
    sections = [s.strip() for s in sections if s.strip()]

    if len(sections) < 3:
        return {
            "purpose": summary_text,
            "company_background": "",
            "submission_details": {"text": "", "bullets": []}
        }

    submission_text = sections[2]
    bullet_lines = re.findall(r"-\s\*\*(.*?)\*\*|-\s+(.*)", submission_text)

    bullets = []
    for b1, b2 in bullet_lines:
        bullet = b1 if b1 else b2
        if bullet and bullet.strip():
            bullets.append(bullet.strip())

    return {
        "purpose": sections[0],
        "company_background": sections[1],
        "submission_details": {
            "text": submission_text,
            "bullets": bullets
        }
    }

def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    if ext == ".pdf":
        reader = PdfReader(file_path)
        text = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
    elif ext == ".docx":
        doc = docx.Document(file_path)
        text = " ".join([p.text for p in doc.paragraphs])
    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    elif ext in [".xls", ".xlsx"]:
            xls = pd.ExcelFile(file_path)
            for sheet in xls.sheet_names:
                df = xls.parse(sheet)
                text += "\n".join(
                    df.fillna("").astype(str).agg(" ".join, axis=1)
                ) + "\n" 
    return text.strip()



async def generate_summary(
    text: str,
    provider: str = "gpt-4o-mini",
) -> str:
    """Generate summary of an RFP using LLM."""
    prompt = f"Summarize the following RFP in 3-5 paragraphs:\n\n{text[:8000]}"
    system_prompt = "You are an RFP summarizer."
    content = await _complete_with_fallback(provider, prompt, system_prompt)
    return content.strip()


def get_short_name(filename: str) -> str:
    """
    Extract short client name from PDF filename.
    Examples:
        'McLean_Hospital_RFP.pdf' -> 'McLean'
        'Acme_Inc_Submission.pdf' -> 'Acme'
        'State_of_California_Dept_Health.pdf' -> 'California'
    """
    name = filename.rsplit('.', 1)[0]

    name = name.replace('_', ' ').replace('-', ' ')
    tokens = name.split()

    noise = {
        "hospital", "inc", "llc", "ltd", "company", "co", "services",
        "dept", "department", "state", "university", "institute",
        "health", "center", "centre", "proposal", "rfp", "submission"
    }

    cleaned = [t for t in tokens if t.lower() not in noise]

    if not cleaned:
        return "the organization"
    return cleaned[0].title()

def bump_version(version: str) -> str:
    if not isinstance(version, str) or not version.strip():
        return None

    parts = version.strip().split(".")

    if not all(p.isdigit() for p in parts):
        return None

    if len(parts) == 1:
        (major,) = map(int, parts)
        return str(major + 1)

    if len(parts) == 2:
        major, minor = map(int, parts)
        return f"{major}.{minor + 1}"

    if len(parts) == 3:
        major, minor, patch = map(int, parts)
        return f"{major}.{minor}.{patch + 1}"

    return None

async def get_next_index(rfp_id: int, user_id: int, question: str, db: AsyncSession) -> str:
    last_question_result = await db.execute(
        select(RFPQuestion.question_text)
        .filter(
            RFPQuestion.rfp_id == rfp_id,
            RFPQuestion.admin_id == user_id
        )
        .order_by(RFPQuestion.id.desc())
    )
    last_question = last_question_result.scalars().first()
    index = last_question.split(" ")[0] if last_question else "0"

    new_index = None
    if index:
        try:
            new_index = bump_version(index)
        except Exception as e:
            new_index = "0"
    
    return f"{new_index} {question}"

def clean_extracted_text(text: str) -> str:
    """
    Clean up common PDF/OCR noise so LLM gets better input.
    - Removes obvious page markers
    - Normalizes line breaks and extra spaces
    """
    text = re.sub(r"Page\s+\d+\s*\n?", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)

    return text.strip()

async def get_active_keystone_text(db: AsyncSession, admin_id: int) -> str:
    keystone_result = await db.execute(
        select(KeystoneFile).filter(KeystoneFile.admin_id == admin_id).order_by(KeystoneFile.uploaded_at.desc())
    )
    keystone = keystone_result.scalars().first()
    if not keystone:
        raise HTTPException(
            status_code=400,
            detail="Keystone Data not uploaded. Please upload Keystone XLS."
        )

    return keystone.extracted_text

def extract_xls_text(file_path: str) -> str:
    sheets = pd.read_excel(file_path, sheet_name=None)
    output = []

    for sheet_name, df in sheets.items():
        output.append(f"\n=== {sheet_name} ===\n")
        for _, row in df.iterrows():
            row_text = " | ".join(
                str(cell) for cell in row if pd.notna(cell)
            )
            if row_text.strip():
                output.append(row_text)

    return "\n".join(output)

def delete_rfp_embeddings(file_id: int):
    try:
        namespace = f"rfp_{file_id}"

        # print("Deleting namespace:", namespace)
        index.delete(delete_all=True, namespace=namespace)
        print(f"Deleted embeddings for {namespace}")

    except Exception as e:
        print(f"Error: {e}")


def classification_QaI(
    rfp_text: str,
    selected_sections: list,
    provider,
    fallback_providers: list[str] = None
) -> dict:
    """
    Classify RFP requirements into Instruction (I), Question (Q), or Both (B) using a 4-step decision tree.
    Returns structured JSON with classification results and summary statistics.
    """
    prompt = classification_prompt(rfp_text, selected_sections)
    # print("Generated classification prompt:", prompt[:50], "...")
    # system_prompt = (
    #     "You are a meticulous RFP requirement classifier. You classify each requirement as Instruction (I), Question (Q), or Both (B) based on a strict 4-step decision tree. "
    #     "You NEVER add information not present in the requirement text. You NEVER misclassify based on assumptions — only the explicit text. "
    #     "You ALWAYS identify the first main verb as the key signal for classification. "
    #     "Your output is a single JSON object with detailed classification results and summary statistics, following the exact structure specified in the prompt."
    # )
    content = _complete_with_fallback(provider, prompt, fallback_providers=fallback_providers)
    print("Raw LLM output for classification:", content)

    try:
        return json.loads(content)
    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail="LLM returned invalid JSON for classification results.")
