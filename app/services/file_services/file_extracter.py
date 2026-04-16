import os
import logging
from pathlib import Path

import docx
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation

logger = logging.getLogger(__name__)

# Supported extensions mapped to their handler
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xls", ".xlsx"}


class ExtractionError(Exception):
    """Raised when text extraction fails for a file."""
    pass


def extract_text_from_file(file_path: str) -> str:
    """
    Extract plain text from a supported file.

    Supports: PDF, DOCX, PPTX, XLS, XLSX.

    Args:
        file_path: Absolute or relative path to the file.

    Returns:
        Extracted text as a single stripped string.

    Raises:
        FileNotFoundError: If the file does not exist.
        ValueError: If the file extension is not supported.
        ExtractionError: If extraction fails due to a corrupt or unreadable file.
    """
    path = Path(file_path).resolve()

    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    if not path.is_file():
        raise ValueError(f"Path is not a file: {path}")

    ext = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file extension '{ext}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    try:
        text = _EXTRACTORS[ext](path)
    except (FileNotFoundError, ValueError):
        raise
    except Exception as e:
        logger.exception("Failed to extract text from '%s'", path)
        raise ExtractionError(f"Could not extract text from '{path}': {e}") from e

    return text.strip()




def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))

    if reader.is_encrypted:
        raise ExtractionError(f"PDF is encrypted and cannot be read: {path}")

    pages: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            page_text = page.extract_text() or ""
            pages.append(page_text)
        except Exception as e:
            logger.warning("Skipping PDF page %d in '%s': %s", i, path, e)

    return "\n".join(pages)


def _extract_docx(path: Path) -> str:
    doc = docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    chunks: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            shape_text = shape.text.strip()
            if shape_text:
                slide_texts.append(shape_text)

        if slide_texts:
            chunks.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

    return "\n\n".join(chunks)


def _extract_excel(path: Path) -> str:
    xls = pd.ExcelFile(str(path))
    sheet_chunks: list[str] = []

    for sheet_name in xls.sheet_names:
        df: pd.DataFrame = xls.parse(sheet_name)

        if df.empty:
            continue

        # Stringify every cell, drop rows that are entirely blank
        rows = (
            df.fillna("")
            .astype(str)
            .replace(r"^\s*$", "", regex=True)
        )
        non_empty_rows = rows[rows.apply(lambda r: r.str.strip().any(), axis=1)]
        row_lines = non_empty_rows.agg(" ".join, axis=1).tolist()

        if row_lines:
            sheet_chunks.append(f"[Sheet: {sheet_name}]\n" + "\n".join(row_lines))

    return "\n\n".join(sheet_chunks)


# Dispatch table — avoids if/elif chains and makes adding formats trivial
_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xls":  _extract_excel,
    ".xlsx": _extract_excel,
}