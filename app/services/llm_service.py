import os
import fitz  # PyMuPDF
from dotenv import load_dotenv
import requests
from sqlalchemy.orm import Session
from passlib.context import CryptContext
from app.models.rfp_models import User
from app.config import client, index,SERPAPI_KEY
from pptx import Presentation
from PyPDF2 import PdfReader
from fastapi import HTTPException
import re
import docx


load_dotenv()
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")

def extract_text_from_pdf(pdf_file: bytes) -> str:
    text = ""
    with fitz.open(stream=pdf_file, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def chat_model(model: str, system_prompt: str, user_prompt: str, temperature: float , max_tokens: int) -> str:
    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=temperature,
        max_tokens=max_tokens
    )
    return response.choices[0].message.content.strip()

def generate_search_queries(rfp_text: str) -> list:
    """
    Generate exactly 12 highly targeted Google search queries based on RFP text.
    """
    user_prompt = f"""
    You are an expert market intelligence researcher specializing in analyzing companies from Request for Proposal (RFP) documents.

    Your task: Generate exactly 12 highly targeted Google search queries based only on the RFP text below.  

    The queries must:
    - Be precise, varied, and investigative.
    - Always incorporate unique identifiers from the RFP (company name, product names, technologies, industries).
    - Cover multiple areas:  
      1. Company history and ownership  
      2. Core products, services, or solutions  
      3. Industry verticals or markets served  
      4. Partnerships, clients, and case studies  
      5. Locations and employee count  
      6. Awards, recognition, and certifications  
      7. Financials, funding, or revenue (if available)  
      8. Competitors and market positioning  
      9. Technology platforms mentioned in the RFP  
      10. Recent press releases or news coverage
      11. Proposal submission requirements for this RFP
      12. Official submission due date / deadlines for this RFP

    Format:
    - Output as a bullet list, one query per line.
    - Do not add explanations — only the search queries.

    RFP Text:
    \"\"\"
    {rfp_text}
    \"\"\"
    """

    system_prompt = "You generate Google search queries to build complete company profiles from RFPs."

    content = chat_model(
        model="gpt-4o-mini",    
        system_prompt=system_prompt,
        user_prompt=user_prompt,
        temperature=0.4,
        max_tokens=1000
    )

    return [line.strip(" -•") for line in content.split("\n") if line.strip()]

def extract_company_background_from_rfp(rfp_text: str) -> str:
    """
    Extracts 3 fully detailed sections from an RFP:
    1. Purpose of the RFP
    2. Company Background
    3. Submission Details & Requirements
    """

    user_prompt = f"""
    You are a senior RFP analyst with deep expertise in procurement, compliance, and government/enterprise documentation.

    Your task is to extract and reorganize information from the provided RFP text into **exactly three sections**.  
    You must pull the content directly from the text without losing any detail.

    ============================================================
     **CRITICAL EXTRACTION RULES (DO NOT VIOLATE)**  
    ============================================================

    **1. NO HALLUCINATIONS.**  
    - If the RFP text does not contain a detail, leave it out.  
    - Never infer names, dates, budgets, processes, or company info.

    **2. DO NOT OMIT ANY INFORMATION.**  
    - If relevant content appears in multiple places in the RFP, gather all of it.  
    - Consolidate it into the correct section.

    **3. USE VERBATIM TEXT WHENEVER POSSIBLE.**  
    - Preserve original wording, formatting style, terminology, and phrasing.  
    - Only merge, compress, or rewrite when necessary for clarity.

    **4. STRICT SECTION SEPARATION.**  
    - No submission details in Section 1 or Section 2.  
    - No background info in Section 1 or Section 3.  
    - No purpose-related narrative in Section 2 or Section 3.

    **5. DO NOT DUPLICATE RFP HEADINGS.**  
    - Use only the section headings defined below.

    ============================================================
     **SECTION REQUIREMENTS**
    ============================================================

    ### **Section 1: Purpose of the RFP**
    Extract all content that explains:
    - The purpose, intent, or reason for issuing the RFP.  
    - Strategic goals, problem statements, motivations, and desired outcomes.  
    - Stakeholders, agencies, departments, or sponsoring bodies explicitly tied to purpose.  
    - Scope elements *only when directly connected to purpose*.  
    **Exclude** submission instructions, dates, addresses, or proposal formatting.

    ### **Section 2: Company Background**
    Extract every piece of organizational context about the issuer, including:
    - Full organization name (legal name, aliases, abbreviations).  
    - Company type (public, private, nonprofit, government, etc.).  
    - Headquarters, office locations, regions served, or jurisdiction.  
    - Mission, vision, mandate, history, values, or strategic priorities.  
    - Size, capacity, funding sources, budgets, staff counts.  
    - Programs, services, lines of business, or operational areas.  
    - Governance, leadership, key stakeholders, partner organizations.  
    - DEI or policy priorities (if present).  
    Collect all background across the entire document, even if scattered.

    **Explicit Exclusions**  
    - NO submission rules  
    - NO proposal requirements  
    - NO selection criteria  
    - NO deadlines or contacts  

    ### **Section 3: Submission Details & Requirements**
    Extract every procedural, compliance, and submission-related detail, such as:
    - Proposal due dates, times, and delivery deadlines.  
    - Submission methods (email, portal, hardcopy, courier, etc.).  
    - Physical or digital submission addresses.  
    - Required formats (PDF, Word, binders, number of copies, etc.).  
    - Mandatory forms, certifications, affidavits, or attachments.  
    - Eligibility rules and compliance requirements.  
    - Evaluation criteria, scoring, selection or award process.  
    - Contacts: names, titles, phone numbers, emails.  
    - Timelines, Q&A policies, pre-bid meetings, vendor requirements.  
    Combine all procedural content into a single comprehensive section.

    ============================================================
     **OUTPUT FORMAT (STRICT)**
    ============================================================

    Section 1: Purpose of the RFP  
    [full extracted content]

    Section 2: Company Background  
    [full extracted content]

    Section 3: Submission Details & Requirements  
    [full extracted content]

    ============================================================
     **SOURCE RFP TEXT**
    ============================================================
    \"\"\"
    {rfp_text}
    \"\"\"
    """

    system_prompt = (
        "You are a senior RFP extraction analyst. You extract purpose, background, and submission "
        "requirements from RFPs with perfect accuracy, zero hallucinations, and strict adherence to sections."
    )

    return chat_model(
        model="gpt-4o-mini",
        system_prompt=system_prompt,
        user_prompt=user_prompt,
        temperature=0.1,
        max_tokens=2200
    )

def summarize_results_with_llm(all_snippets: list, rfp_company_text: str) -> str:
    """
    Combine RFP company description and web search snippets into a
    structured, executive-level analysis with 3 fixed sections.
    Ensures Submission Details & Requirements are comprehensive,
    accurate, and formatted as a bullet list with clean spacing.
    """

    combined_snippets = "\n".join(all_snippets)

    user_prompt = f"""
    You are a senior strategy consultant preparing a formal RFP analysis brief.
    
    Critical Formatting & Content Rules:
    - Produce exactly three sections in the order specified: Purpose of the RFP, Company Background, and Submission Details & Requirements.
    - Insert a blank line between paragraphs in Sections 1 and 2 for readability.
    - In Section 3, extract **all** submission-related details and requirements from the RFP, even if scattered across multiple sections or repeated.
    - Present Section 3 as a bullet list, with each requirement or detail on its own line, quoted or restated verbatim to preserve the original wording and intent.
    - Do not paraphrase, summarize, or omit any submission requirements in Section 3.
    - Exclude submission deadlines, contact details, or procedural instructions from Section 1 and Section 2.
    - Use information from web snippets to enhance Section 2 (Company Background) only if it is verified, relevant, and complements the RFP content.
    - If information is missing for any section, include a note stating: "No relevant information provided in the RFP or web snippets."
    - If conflicting information exists (e.g., between RFP and snippets), prioritize RFP data and note discrepancies in Section 2 (e.g., "Web snippets suggest [X], but RFP states [Y]").
    - Ensure the output is professional, concise, and avoids redundancy while maintaining all required details.
    
    Do Not:
    - Include submission deadlines, contact details, or procedural instructions in Section 1 or Section 2.
    - Paraphrase or modify submission requirements in Section 3; use exact wording or faithful restatements.
    - Introduce speculative or unverified information not present in the RFP or web snippets.
    - Use Markdown or other formatting (e.g., bold, asterisks) in the output; use plain text with the exact section headings provided.
    - Repeat section headings within the extracted content.
    
    ---
    **Section 1: Purpose of the RFP**
    [Full paragraph explanation of why the RFP was issued, its goals, scope, and strategic drivers.
    Exclude submission deadlines and contacts here.]

    **Section 2: Company Background**
    [Full paragraph company profile combining RFP content and verified details from the web.
    Include: company name, founding year, HQ, ownership, core offerings, markets served,
    strategic initiatives, awards, major clients/partners, and market position.]

    **Section 3: Submission Details & Requirements**
    [Bullet list of every requirement and operational detail from the RFP,
    including submission due date, question deadline, contact names/emails,
    submission method, required proposal contents, eligibility criteria,
    and any special instructions or conditions.]
    ---

    RFP Company Description:
    \"\"\"
    {rfp_company_text}
    \"\"\"

    Web Search Snippets:
    \"\"\"
    {combined_snippets}
    \"\"\"
    """

    system_prompt = (
        "You produce structured three-part RFP summaries. "
        "Always include every submission detail in Section 3, formatted as a bullet list. "
        "Maintain clear formatting with line breaks after paragraphs."
    )

    return chat_model(
        model="gpt-4o-mini",
        system_prompt=system_prompt,
        user_prompt=user_prompt,
        temperature=0.2,  
        max_tokens=2200 
    )

def search_with_serpapi(query: str):
    url = "https://serpapi.com/search"
    params = {
        "engine": "google",
        "q": query,
        "api_key": SERPAPI_KEY
    }
    res = requests.get(url, params=params)
    data = res.json()

    results = []
    if "organic_results" in data:
        for item in data["organic_results"][:5]:
            results.append({
                "title": item.get("title"),
                "link": item.get("link"),
                "snippet": item.get("snippet")
            })
    return results

def extract_questions_with_llm(pdf_text: str) -> dict:
    prompt = f"""
        You are a **Senior RFP Analysis Expert** with extensive experience in analyzing government and corporate procurement documents. Your task is to **extract every explicit or implied vendor response requirement** from the provided RFP text with utmost accuracy.

        ---

        ###  Objective:
        Extract **every question, instruction, or directive** that requires a vendor to provide information, documentation, confirmation, or explanation in their proposal.

        ---

        ###  Do NOT:
        - Do **not** summarize, paraphrase, or alter the wording of any question or instruction.
        - Do **not** infer or invent questions not explicitly or implicitly requiring a vendor response.
        - Do **not** include:
        - Administrative or procedural details (e.g., submission dates, formats, or contact info).
        - References or appendices instructions (e.g., “Include three client references”).
        - Informational context that does **not** request a vendor action (e.g., “The agency seeks to improve outreach”).
        - Background, introduction, or goal statements unless they **explicitly** ask for a response.
        - Rhetorical or explanatory questions that are not vendor-facing (e.g., “Why is this important?”).

        ---

        ###  Do:
        - Extract **verbatim** text for each question or instruction, including punctuation and capitalization exactly as in the RFP.
        - Treat **any directive phrased as a statement that clearly expects a vendor response** as a question.  
        Examples:
        - “The vendor shall provide…” → extract as a response-required instruction.  
        - “Offeror must demonstrate…” → extract as a required response.  
        - Preserve **section hierarchy, numbering, and structure** from the source RFP.
        - Use and extend the numbering structure logically:
        - “1.2 Proposal Requirements” → “1.2.1”, “1.2.2”, etc.
        - “III.A Technical Approach” → “III.A.1”, “III.A.2”, etc.
        - “I Introduction” → “I.1”, “I.2”, etc.
        - Restart numbering at **1** within each new section.
        - Group extracted questions **under their exact section headings** and **include the original section number/title**.
        - If a section includes **no vendor-response elements**, **omit** it from the output.
        - If the document lacks formal numbering or structure, group under **“General Questions” (G.1, G.2, etc.)**.

        ---

        ###  Edge Case Handling:
        - If mixed content appears (narrative + questions), include **only the explicit vendor-facing directives**.
        - Recognize both **direct** (“Describe your process”) and **implicit** (“Vendor shall provide…”) vendor response requirements.
        - Handle nested numbering or outline styles accurately (e.g., “1.2.3.1” or “III.B.1.a”).
        - For bullet points, lists, or tables, extract each question/instruction as a **single clean text string**, ignoring formatting unless it changes meaning.
        - Preserve **logical hierarchy** even when numbering is inconsistent or missing by analyzing indentation, font cues (if provided), or section headers.

        ---

        ###  Output Format:
        Return a **strict JSON object** (no Markdown, no commentary, no explanations).

        Each section in the JSON must contain:
        - `"section"` → exact section title and number as written in the RFP.
        - `"questions"` → list of extracted, verbatim vendor-response items prefixed with their section-based numbering.

        ---

        ###  Example Output:
        {{
        "1": {{
            "section": "1.1 Scope of Work",
            "questions": [
            "1.1.1 Describe your agency's approach to developing a comprehensive media strategy.",
            "1.1.2 Provide a sample timeline for implementation."
            ]
        }},
        "2": {{
            "section": "2.1 Proposal Requirements",
            "questions": [
            "2.1.1 Outline your firm's qualifications and experience.",
            "2.1.2 Explain how your solution addresses the stated objectives."
            ]
        }},
        "3": {{
            "section": "III.A Technical Approach",
            "questions": [
            "III.A.1 Describe your proposed system architecture.",
            "III.A.2 Provide your data protection and cybersecurity protocols."
            ]
        }}
        }}

        ---

        ###  Input:
        RFP Document Text:
        \"\"\"
        {pdf_text}
        \"\"\"
        """

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": (
                    "You are an expert assistant trained to extract structured, "
                    "section-wise response prompts from RFP documents. "
                    "Always return clean JSON with no extra formatting. "
                    "Never modify, summarize, or rephrase the RFP text."
                ),
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.0,  
        max_tokens=2000,
    )

    content = response.choices[0].message.content

    try:
        import json
        grouped_questions = json.loads(content)
    except Exception:
        grouped_questions = {"raw_text": content}

    return grouped_questions

def get_user_by_email(db: Session, email: str):
    return db.query(User).filter(User.email == email).first()

def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def hash_password(password):
    return pwd_context.hash(password)

def get_similar_context(question: str, rfp_id: int, top_k: int = 5):
    """
    Retrieve RFP-specific chunks from Pinecone using file_id metadata filter.
    """
    try:
        embedding = client.embeddings.create(
            input=[question],
            model="text-embedding-3-small"
        ).data[0].embedding

        results = index.query(
            vector=embedding,
            top_k=top_k,
            include_metadata=True,
            filter={"file_id": str(rfp_id)}
        )

        context_texts = [match["metadata"]["text"] for match in results["matches"]]
        sources = [
            {
                "score": match["score"],
                "file_id": match["metadata"].get("file_id"),
                "chunk_index": match["metadata"].get("chunk_index"),
                "snippet": match["metadata"].get("text")[:300],
            }
            for match in results["matches"]
        ]

        return "\n".join(context_texts), sources

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Pinecone retrieval failed: {str(e)}")

def generate_answer_with_context(question: str, context: str, short_name: str) -> str:
    prompt = f"""
    You are an expert proposal writer. You write RFP responses **only using the provided context**.

    ### Client Name Rule (MANDATORY)
    - Refer to the issuer exclusively as "{short_name}".
    - Never use “the client”.
    - Never guess or invent a different name.

    ### Voice & Point of View (MANDATORY)
    - Use “we” or “our”.
    - Never use “I”.
    - Never refer to Ringer in third person.

    ### Pricing Rules (MANDATORY)
    - Ringer uses flat-rate pricing only.
    - Never mention hourly rates, per-hour billing, or hourly estimates.
    - Do NOT generate or assume any hourly cost structure.
    - If the context mentions hourly pricing, rewrite it into a flat-rate model without creating specific prices.

    ### Subcontractor / Vendor Operating Model (MANDATORY)
    - Ringer delivers services through **subcontractors and external vendors**.
    - This is a core part of how Ringer operates.
    - When describing service delivery, staffing, capabilities, or processes:
        - reflect that subcontractors/vendors are used.
        - never imply Ringer relies solely on in-house full-time staff.
    - Keep descriptions accurate and high-level unless the context specifies otherwise.

    ### Accuracy Rules
    - Use ONLY the information in the context.
    - If the context lacks the information, write:
      “We do not have enough information to provide that detail based on the available context.”

    ### Tone & Style
    - Professional, concise, confident.
    - No vague marketing language.

    ### Formatting
    - No bullet points unless context explicitly requires it.
    - Do NOT mention “context” or “question”.

    ----

    Context:
    {context}

    Question:
    {question}

    Final Answer:
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a highly skilled RFP response specialist who strictly follows instructions."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.2,
            max_tokens=1000
        )
        return response.choices[0].message.content.strip()

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM generation failed: {str(e)}")

def analyze_answer_score_only(question_text: str, answer_text: str) -> float:
    prompt = f"""
You are acting as a strict RFP evaluator.

Your task:
- Carefully read the question and the provided answer.
- Judge how well the answer directly addresses the question.
- Give a score from 0.0 to 10.0 based on relevance, clarity, and completeness.

Scoring Rules:
- 0.0 → No relevance or completely incorrect
- 1.0-3.0 → Very poor / barely addresses the question
- 4.0-6.0 → Partially addresses the question, with gaps
- 7.0-8.5 → Good, mostly complete but could be stronger
- 9.0-10.0 → Excellent, fully relevant and comprehensive

 Output Instruction:
Return ONLY the numeric score as a float (e.g., `7.5`, `9.0`, `0.0`). 
Do NOT include words, labels, or explanations.

Question: {question_text}

Answer: {answer_text}
    """

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2,
    )

    score_text = response.choices[0].message.content.strip()
    try:
        return float(score_text)
    except ValueError:
        return None

def parse_rfp_summary(summary_text: str):
    """
    Convert LLM markdown summary into structured JSON:
    {
      "purpose": "...",
      "company_background": "...",
      "submission_details": {
          "text": "...",
          "bullets": ["...", "...", ...]
      }
    }
    """
    summary_text = re.sub(r"^#+\s*RFP Analysis Brief\s*", "", summary_text, flags=re.IGNORECASE)

    sections = re.split(r"####\s*Section\s*\d:\s*", summary_text)
    sections = [s.strip() for s in sections if s.strip()]

    if len(sections) < 3:
        return {
            "purpose": summary_text,
            "company_background": "",
            "submission_details": {"text": "", "bullets": []}
        }

    submission_text = sections[2]
    bullet_lines = re.findall(r"-\s\*\*(.*?)\*\*|-\s+(.*)", submission_text)

    bullets = []
    for b1, b2 in bullet_lines:
        bullet = b1 if b1 else b2
        if bullet and bullet.strip():
            bullets.append(bullet.strip())

    return {
        "purpose": sections[0],
        "company_background": sections[1],
        "submission_details": {
            "text": submission_text,
            "bullets": bullets
        }
    }

def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    if ext == ".pdf":
        reader = PdfReader(file_path)
        text = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
    elif ext == ".docx":
        doc = docx.Document(file_path)
        text = " ".join([p.text for p in doc.paragraphs])
    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    return text.strip()

def get_embedding(text: str):
    response = client.embeddings.create(
        model="text-embedding-3-small",
        input=text
    )
    return response.data[0].embedding

def generate_summary(text: str) -> str:
    """Generate summary of an RFP using LLM."""
    prompt = f"Summarize the following RFP in 3-5 paragraphs:\n\n{text[:8000]}"
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": "You are an RFP summarizer."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content.strip()

def delete_rfp_embeddings(file_id: int):
    try:
        results = index.query(
            vector=[0.0] * 1536,
            top_k=10000,   
            include_metadata=False,
            filter={"file_id": str(file_id)}
        )

        vector_ids = [match["id"] for match in results.get("matches", [])]
        if vector_ids:
            index.delete(ids=vector_ids)
            print(f" Deleted {len(vector_ids)} Pinecone vectors for file_id {file_id}")
        else:
            print(f" No Pinecone vectors found for file_id {file_id}")

    except Exception as e:
        print(f" Error deleting embeddings for RFP {file_id}: {e}")

def get_short_name(filename: str) -> str:
    """
    Extract short client name from PDF filename.
    Examples:
        'McLean_Hospital_RFP.pdf' -> 'McLean'
        'Acme_Inc_Submission.pdf' -> 'Acme'
        'State_of_California_Dept_Health.pdf' -> 'California'
    """
    name = filename.rsplit('.', 1)[0]

    name = name.replace('_', ' ').replace('-', ' ')

    tokens = name.split()

    noise = {
        "hospital", "inc", "llc", "ltd", "company", "co", "services",
        "dept", "department", "state", "university", "institute",
        "health", "center", "centre", "proposal", "rfp", "submission"
    }

    cleaned = [t for t in tokens if t.lower() not in noise]

    if not cleaned:
        return "the organization"
    return cleaned[0].title()






